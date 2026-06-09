import os
from pathlib import Path
import fitz  # PyMuPDF
import docx
from pptx import Presentation
def validate_file_path(file_path: str) -> Path:
    path = Path(file_path).resolve()
    base = Path(__file__).parent.resolve()

    allowed_dirs = [
        base / "uploads",
        base / "vector_store",
    ]

    for allowed_dir in allowed_dirs:
        try:
            path.relative_to(allowed_dir)
            return path
        except ValueError:
            continue

    raise ValueError(f"Access denied: file path outside allowed directories")

def extract_text(file_path: str) -> str:
    safe_path = validate_file_path(file_path)
    ext = safe_path.suffix.lower()

    if ext == ".pdf":
        return extract_pdf_text(str(safe_path))
    elif ext == ".docx":
        return extract_docx_text(str(safe_path))
    elif ext == ".pptx":
        return extract_pptx_text(str(safe_path))
    elif ext == ".txt":
        return extract_txt(str(safe_path))
    else:
        raise ValueError(f"Unsupported file type: {ext}")



def extract_pdf_text(file_path: str) -> str:
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
    return text.strip()

def extract_docx_text(file_path: str) -> str:
    doc = docx.Document(file_path)
    full_text = []
    
    # Extract text from paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            full_text.append(para.text)
    
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    full_text.append(cell.text)
    
    return "\n".join(full_text).strip()

def extract_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_txt(file_path: str) -> str:
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read().strip()


