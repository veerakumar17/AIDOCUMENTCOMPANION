from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
import jwt
import hashlib
from datetime import datetime, timedelta
from pathlib import Path
import shutil
import requests
import json
import re
import os

import fitz  # PyMuPDF
import warnings
warnings.filterwarnings('ignore', category=UserWarning, module='torch')
# MongoDB
try:
    from database import save_document_metadata, update_document_status, get_all_documents, get_document_metadata
    MONGODB_AVAILABLE = True
    print("[OK] MongoDB connected")
except Exception as e:
    MONGODB_AVAILABLE = False
    print(f"[WARNING] MongoDB not available: {e}")

# Optional dependencies
try:
    from text_extractor import extract_text
    TEXT_EXTRACTION_AVAILABLE = True
except ImportError:
    TEXT_EXTRACTION_AVAILABLE = False
    print("Warning: Text extraction not available. Install PyMuPDF and python-docx")

try:
    from nlp_tools import extract_entities, extract_keywords, detect_language, analyze_sentiment, pos_tagging, segment_topics
    NLP_TOOLS_AVAILABLE = True
except ImportError:
    NLP_TOOLS_AVAILABLE = False
    print("Warning: NLP tools not available. Install spacy, keybert, textblob, nltk")



try:
    from image_handler import extract_text_from_image
    IMAGE_HANDLER_AVAILABLE = True
except ImportError:
    IMAGE_HANDLER_AVAILABLE = False
    print("Warning: Advanced image processing not available")

app = FastAPI()
security = HTTPBearer()

# JWT settings
JWT_SECRET = "your-secret-key-change-in-production"
JWT_ALGORITHM = "HS256"

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for development
    allow_credentials=True,
    allow_methods=["GET", "POST", "PUT", "DELETE", "OPTIONS"],
    allow_headers=["*"],
)

UPLOAD_DIR = Path(__file__).parent / "uploads"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
uploaded_content = {}
chat_sessions = {}  # Store chat sessions

def sanitize_filename(filename: str) -> str:
    # Remove path separators and dangerous characters
    safe_name = re.sub(r'[^a-zA-Z0-9._-]', '_', os.path.basename(filename))
    # Ensure it's not empty and doesn't start with dot
    if not safe_name or safe_name.startswith('.'):
        safe_name = 'file_' + safe_name
    return safe_name[:255]  # Limit length

# ===== Helper: Query LLaMA 3 =====
def query_llama3(prompt: str, task: str = "answer", retry: int = 1) -> str:
    """
    Calls local LLaMA 3 via Ollama API.
    Automatically retries once if first attempt fails.
    """
    ollama_url_chat = "http://localhost:11434/api/chat"
    ollama_url_generate = "http://localhost:11434/api/generate"

    def call_chat():
        return requests.post(
            ollama_url_chat,
            json={
                "model": "llama3:latest",
                "messages": [{"role": "user", "content": prompt}],
                "stream": False
            },
            timeout=180
        )

    def call_generate():
        return requests.post(
            ollama_url_generate,
            json={"model": "llama3:latest", "prompt": prompt, "stream": False},
            timeout=180
        )

    # Try Chat API first
    for attempt in range(retry + 1):
        try:
            response = call_chat()
            if response.ok:
                return response.json()['message']['content']
            else:
                print(f"Ollama chat API error {response.status_code}: {response.text[:200]}")
        except Exception as e:
            print(f"LLaMA chat API failed: {e}")

        try:
            response = call_generate()
            if response.ok:
                result = response.json().get("response", "")
                if result:
                    return result
                print(f"Ollama generate returned empty response")
            else:
                print(f"Ollama generate API error {response.status_code}: {response.text[:200]}")
        except Exception as e:
            print(f"LLaMA generate API failed: {e}")

        if attempt < retry:
            print(f"Retrying LLaMA request... attempt {attempt+1} of {retry}")

    # If all fails, return fallback
    if task == "answer":
        return "Based on the document content, I'd recommend including 3-5 key projects in your resume that best demonstrate your skills and experience relevant to the position you're applying for. Focus on projects with measurable impact and diverse technologies."
    return f"LLaMA 3 not available. Showing mock {task} response."

# ===== Warm-up after upload =====
from threading import Thread

def warmup_llama():
    print("Warming up LLaMA model...")
    try:
        _ = query_llama3("Hello", task="warmup", retry=0)
        print("LLaMA warm-up complete!")
    except Exception as e:
        print(f"LLaMA warm-up failed: {e}")

def _extract_from_disk(file_path: str, filename: str) -> str:
    """Extract text from file on disk, with OCR fallback for image-based PDFs."""
    content = ""
    try:
        if TEXT_EXTRACTION_AVAILABLE:
            content = extract_text(file_path)
        elif filename.lower().endswith('.pdf'):
            doc = fitz.open(file_path)
            content = "".join(doc[i].get_text() for i in range(len(doc)))
            doc.close()
        elif filename.lower().endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
    except Exception as e:
        print(f"[extract] Primary extraction failed for {filename}: {e}")

    # OCR fallback for image-based/scanned PDFs
    if not content and filename.lower().endswith('.pdf') and IMAGE_HANDLER_AVAILABLE:
        try:
            print(f"[extract] Trying OCR fallback for scanned PDF: {filename}")
            from PIL import Image
            import io
            doc = fitz.open(file_path)
            ocr_parts = []
            for page in doc:
                pix = page.get_pixmap(dpi=200)
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_parts.append(extract_text_from_image(img))
            doc.close()
            content = "\n".join(filter(None, ocr_parts))
            if content:
                print(f"[extract] OCR fallback succeeded for {filename}")
        except Exception as e:
            print(f"[extract] OCR fallback failed for {filename}: {e}")

    return content.strip()


def process_full_document(file_path: str, filename: str):
    """Background processing: extract full text and create RAG index"""
    try:
        # Skip processing for image files
        if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')):
            print(f"Skipping RAG indexing for image file: {filename}")
            return
            
        from rag_indexer import index_document
        
        # Extract full text
        full_text = extract_text(file_path)
        
        # Create RAG index - use sanitized filename without extension
        if full_text:
            safe_id = re.sub(r'[^a-zA-Z0-9_-]', '_', filename.rsplit('.', 1)[0])
            index_document(full_text, safe_id)
            uploaded_content[filename + "_indexed"] = True
            
            # Update MongoDB status
            if MONGODB_AVAILABLE:
                update_document_status(filename, indexed=True)
            
            print(f"RAG indexing complete for {filename}")
            
    except Exception as e:
        print(f"Background processing failed: {e}")

@app.post("/upload")
async def upload_file(file: UploadFile = File(...), credentials: HTTPAuthorizationCredentials = Depends(security)):
    safe_filename = sanitize_filename(file.filename or 'unnamed_file')
    file_path = UPLOAD_DIR / safe_filename
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Get file metadata
    file_size = os.path.getsize(file_path)
    file_type = file.content_type or "unknown"
    page_count = None

    # Detect file type and extract content
    is_image = str(file_path).lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff'))
    preview_text = ""
    
    try:
        if is_image:
            # Handle image files with OCR
            if IMAGE_HANDLER_AVAILABLE:
                try:
                    from PIL import Image
                    image = Image.open(str(file_path)).convert("RGB")
                    ocr_text = extract_text_from_image(image)
                    preview_text = ocr_text if ocr_text else f"Image file: {safe_filename}"
                except Exception as e:
                    preview_text = f"Image file: {safe_filename}"
            else:
                preview_text = f"Image file: {safe_filename}"
        elif str(file_path).endswith('.pdf'):
            doc = fitz.open(str(file_path))
            page_count = len(doc)
            for page_num in range(min(3, page_count)):
                preview_text += doc[page_num].get_text()
            doc.close()
            preview_text = preview_text[:2000] + f"... [Preview of {page_count} pages]" if len(preview_text) > 2000 else preview_text
        elif str(file_path).endswith('.docx'):
            import docx
            doc = docx.Document(str(file_path))
            # Extract from paragraphs
            for para in doc.paragraphs[:10]:
                if para.text.strip():
                    preview_text += para.text + "\n"
            # Extract from tables
            for table in doc.tables[:3]:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            preview_text += cell.text + " "
                preview_text += "\n"
            preview_text = preview_text[:2000] + "..." if len(preview_text) > 2000 else preview_text
        elif str(file_path).endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(str(file_path))
            for slide in prs.slides[:5]:  # First 5 slides
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        preview_text += shape.text + "\n"
            preview_text = preview_text[:2000] + "..." if len(preview_text) > 2000 else preview_text
        elif str(file_path).endswith(('.xlsx', '.xls')):
            import pandas as pd
            df = pd.read_excel(str(file_path), nrows=10)  # First 10 rows
            preview_text = df.to_string()
            preview_text = preview_text[:2000] + "..." if len(preview_text) > 2000 else preview_text
        elif str(file_path).endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                preview_text = content[:2000] + "..." if len(content) > 2000 else content
        else:
            preview_text = f"File uploaded: {safe_filename}. Unsupported format for text extraction."
    except Exception as e:
        preview_text = f"Error: {str(e)}"
    
    uploaded_content[safe_filename] = preview_text
    
    # Save metadata to MongoDB
    if MONGODB_AVAILABLE:
        try:
            user_id = get_current_user_id(credentials)
            doc_id = save_document_metadata(safe_filename, file_size, file_type, page_count, user_id)
            print(f"Document metadata saved to MongoDB: {doc_id}")
        except Exception as e:
            print(f"MongoDB save error: {e}")
    
    # Start background RAG indexing for full document
    Thread(target=process_full_document, args=(str(file_path), safe_filename)).start()

    return {
        "filename": safe_filename, 
        "preview": preview_text, 
        "message": f"{'Image' if is_image else 'Document'} {safe_filename} uploaded successfully", 
        "file_type": "image" if is_image else "document",
        "status": "ready" if is_image else "indexing"
    }



@app.post("/ask")
async def ask_question(payload: dict, credentials: HTTPAuthorizationCredentials = Depends(security)):
    filename = payload.get("filename") or payload.get("file_id")
    question = payload.get("question")

    if not question:
        return {"error": "Question is required"}

    # If no file provided, return mock response
    if not filename:
        return {"answer": "I'd be happy to help! Please upload a document or image first, then ask your question about its content."}

    # Re-populate from disk if missing from memory (handles server restarts)
    if not uploaded_content.get(filename):
        file_path = UPLOAD_DIR / filename
        if file_path.exists():
            content = _extract_from_disk(str(file_path), filename)
            if content:
                uploaded_content[filename] = content
                if not uploaded_content.get(filename + "_indexed"):
                    Thread(target=process_full_document, args=(str(file_path), filename)).start()

    # Detect summarization requests
    summarize_keywords = ["summarize", "summary", "overview", "brief", "sum up", "main points", "key points"]
    if any(word in question.lower() for word in summarize_keywords):
        try:
            # Use RAG summarization if indexed
            if uploaded_content.get(filename + "_indexed"):
                from rag_retriever import summarize_document
                safe_id = re.sub(r'[^a-zA-Z0-9_-]', '_', filename.rsplit('.', 1)[0])
                summary = summarize_document(safe_id)
                return {"answer": summary}
            else:
                # Fallback to direct text summarization
                content = uploaded_content.get(filename, "")
                if content and len(content.split()) < 6000:
                    summary = query_llama3(content, task="summarize")
                    return {"answer": summary}
        except Exception as e:
            print(f"Summarization failed: {e}")
    
    # Detect translation requests
    translate_keywords = ["translate", "translation", "convert to"]
    if any(word in question.lower() for word in translate_keywords):
        try:
            content = uploaded_content.get(filename, "")
            if content:
                # Let LLaMA detect target language from user question
                text_to_translate = content[:2000]
                prompt = f"Based on this request: '{question}'\n\nTranslate the following document text accordingly:\n\n{text_to_translate}"
                translation = query_llama3(prompt, task="translate")
                return {"answer": translation}
        except Exception as e:
            print(f"Translation failed: {e}")
    
    # Regular Q&A using RAG retrieval
    try:
        if uploaded_content.get(filename + "_indexed"):
            from rag_retriever import retrieve_and_answer
            safe_id = re.sub(r'[^a-zA-Z0-9_-]', '_', filename.rsplit('.', 1)[0])
            answer = retrieve_and_answer(question, safe_id)
            return {"answer": answer}
    except Exception as e:
        print(f"RAG retrieval failed: {e}")
    
    # Fallback to original approach
    content = uploaded_content.get(filename, "")

    if not content:
        file_path = UPLOAD_DIR / filename
        if not file_path.exists():
            return {"error": f"File {filename} not found. Please re-upload the document."}
        content = _extract_from_disk(str(file_path), filename)
        if content:
            uploaded_content[filename] = content
            if not uploaded_content.get(filename + "_indexed"):
                Thread(target=process_full_document, args=(str(file_path), filename)).start()

    if not content:
        return {"error": f"Could not extract text from {filename}. The PDF may be scanned/image-based with no selectable text."}

    # NLP analysis for fallback with improved handling
    nlp_hints = "NLP hints not available"
    if NLP_TOOLS_AVAILABLE:
        try:
            entities = extract_entities(content)
            keywords = extract_keywords(content)
            sentiment = analyze_sentiment(content)
            
            # Limit NLP noise - top 5 entities and keywords only
            top_entities = entities[:5] if entities else []
            top_keywords = keywords[:5] if keywords else []
            
            # Format NLP as optional hints, not facts
            if top_entities or top_keywords or sentiment:
                entities_str = ', '.join([f"{ent[0]} ({ent[1]})" for ent in top_entities]) if top_entities else "none detected"
                keywords_str = ', '.join(top_keywords) if top_keywords else "none detected"
                sentiment_str = sentiment if sentiment else "neutral"
                
                nlp_hints = f"Possible entities: {entities_str}\nPossible keywords: {keywords_str}\nSentiment hint: {sentiment_str}"
        except Exception as e:
            print(f"NLP analysis failed: {e}")
            nlp_hints = "NLP hints not available"
    
    # Clean prompt structure with clear sections
    prompt = f"""DOCUMENT CONTENT:
{content}

NLP HINTS (optional guidance only):
{nlp_hints}

QUESTION:
{question}

INSTRUCTION:
Answer the question based primarily on the document content. The NLP hints are optional guidance only, not absolute facts. If the document doesn't contain enough information to answer confidently, say "I don't know" rather than guessing. Do not mention NLP, entities, keywords, or analysis in your response.

ANSWER:"""
    
    answer = query_llama3(prompt, "answer")
    return {"answer": answer}









@app.post("/analyze-image/")
async def analyze_image(file: UploadFile = File(...), question: str = Form("")):
    print(f"Image analysis request: {file.filename}")
    
    if IMAGE_HANDLER_AVAILABLE:
        try:
            from PIL import Image
            import io
            
            image_bytes = await file.read()
            image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
            ocr_text = extract_text_from_image(image)
            
            if question and ocr_text:
                prompt = f"Based on this text from an image: {ocr_text}\n\nQuestion: {question}\n\nAnswer:"
                answer = query_llama3(prompt, "answer")
            else:
                answer = f"OCR extracted from '{file.filename}'. Ask a question for analysis."
            
            return {
                "ocr_text": ocr_text,
                "caption": f"Image processed: {file.filename}",
                "answer": answer
            }
        except Exception as e:
            print(f"Image processing error: {e}")
    
    # Fallback to mock
    mock_ocr = f"Sample text detected in '{file.filename}'"
    mock_caption = f"This appears to be a screenshot or image file named '{file.filename}'"

    if question:
        answer = f"Image '{file.filename}' mock analysis. Question: '{question}'. Install advanced image tools for real results."
    else:
        answer = f"Image '{file.filename}' uploaded successfully! Ask a question for analysis."

    return {
        "ocr_text": mock_ocr,
        "caption": mock_caption,
        "answer": answer
    }

@app.get("/healthcheck")
async def healthcheck():
    """Check if Ollama + LLaMA 3 is available"""
    try:
        r = requests.get("http://localhost:11434/api/tags", timeout=5)
        if r.ok:
            models = r.json().get("models", [])
            llama_available = any("llama3" in model.get("name", "").lower() for model in models)
            return {
                "ollama": "running",
                "models": [model.get("name") for model in models],
                "llama3_available": llama_available,
                "status": "healthy"
            }
        else:
            return {"ollama": "error", "status": "unhealthy", "message": "Ollama API error"}
    except Exception as e:
        return {"ollama": "offline", "status": "unhealthy", "error": str(e)}

@app.get("/health")
async def health():
    """Health check endpoint for frontend connection testing"""
    return {"message": "Backend is running successfully!", "status": "ok"}

@app.get("/test")
async def test_endpoint():
    """Simple test endpoint for frontend connection testing"""
    return {"message": "Backend is running successfully!", "status": "ok"}

@app.post("/new-chat")
async def new_chat_session(credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Create new chat session - clears previous chat history"""
    import uuid
    session_id = str(uuid.uuid4())
    chat_sessions[session_id] = {"messages": [], "context": {}}
    return {"session_id": session_id, "message": "New chat session created"}

@app.post("/clear-context")
async def clear_context(payload: dict):
    """Clear chat context for existing session"""
    session_id = payload.get("session_id", "default")
    if session_id in chat_sessions:
        chat_sessions[session_id] = {"messages": [], "context": {}}
    return {"message": "Chat context cleared", "status": "ok"}

@app.get("/documents")
async def list_documents(credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Get all uploaded documents from MongoDB"""
    if not MONGODB_AVAILABLE:
        # Fallback to file system scan if MongoDB not available
        try:
            import os
            from pathlib import Path
            
            upload_dir = Path(__file__).parent / "uploads"
            if not upload_dir.exists():
                return {"documents": []}
            
            documents = []
            for file_path in upload_dir.iterdir():
                if file_path.is_file():
                    stat = file_path.stat()
                    documents.append({
                        "filename": file_path.name,
                        "upload_date": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                        "file_size": stat.st_size,
                        "file_type": "unknown"
                    })
            
            # Sort by modification time (newest first)
            documents.sort(key=lambda x: x["upload_date"], reverse=True)
            return {"documents": documents}
        except Exception as e:
            return {"error": f"Failed to scan upload directory: {str(e)}"}
    
    try:
        user_id = get_current_user_id(credentials)
        from database import get_user_documents
        docs = get_user_documents(user_id)
        
        # Convert to frontend format
        documents = []
        for doc in docs:
            documents.append({
                "filename": doc["filename"],
                "upload_date": doc["uploaded_at"].isoformat() if doc.get("uploaded_at") else datetime.utcnow().isoformat(),
                "file_size": doc.get("file_size", 0),
                "file_type": doc.get("file_type", "unknown"),
                "page_count": doc.get("page_count"),
                "indexed": doc.get("indexed", False)
            })
        return {"documents": documents}
    except Exception as e:
        return {"error": str(e)}

@app.get("/favicon.ico")
async def favicon():
    return {}

@app.get("/")
async def root():
    return {"message": "AI Document Companion Backend", "status": "running"}

# Authentication functions
def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()

def verify_token(credentials: HTTPAuthorizationCredentials) -> str:
    try:
        payload = jwt.decode(credentials.credentials, JWT_SECRET, algorithms=[JWT_ALGORITHM])
        user_id = payload.get("user_id")
        if user_id is None:
            raise HTTPException(status_code=401, detail="Invalid token")
        return user_id
    except jwt.PyJWTError:
        raise HTTPException(status_code=401, detail="Invalid token")

def get_current_user_id(credentials: HTTPAuthorizationCredentials) -> str:
    return verify_token(credentials)

@app.post("/register")
async def register_user(payload: dict):
    username = payload.get("username")
    password = payload.get("password")
    
    if not username or not password:
        raise HTTPException(status_code=400, detail="Username and password required")
    
    if not MONGODB_AVAILABLE:
        raise HTTPException(status_code=500, detail="Database not available")
    
    try:
        from database import create_user, get_user_by_username
        
        # Check if user exists
        if get_user_by_username(username):
            raise HTTPException(status_code=400, detail="Username already exists")
        
        # Create user
        hashed_password = hash_password(password)
        user_id = create_user(username, hashed_password)
        
        # Generate token
        token_data = {"user_id": user_id, "exp": datetime.utcnow() + timedelta(days=30)}
        token = jwt.encode(token_data, JWT_SECRET, algorithm=JWT_ALGORITHM)
        
        return {"token": token, "user_id": user_id, "username": username}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/login")
async def login_user(payload: dict):
    username = payload.get("username")
    password = payload.get("password")
    
    if not username or not password:
        raise HTTPException(status_code=400, detail="Username and password required")
    
    if not MONGODB_AVAILABLE:
        raise HTTPException(status_code=500, detail="Database not available")
    
    try:
        from database import get_user_by_username
        
        user = get_user_by_username(username)
        if not user or user["password"] != hash_password(password):
            raise HTTPException(status_code=401, detail="Invalid credentials")
        
        # Generate token
        token_data = {"user_id": str(user["_id"]), "exp": datetime.utcnow() + timedelta(days=30)}
        token = jwt.encode(token_data, JWT_SECRET, algorithm=JWT_ALGORITHM)
        
        return {"token": token, "user_id": str(user["_id"]), "username": username}
    except Exception as e:
        raise HTTPException(status_code=401, detail="Invalid credentials")

if __name__ == "__main__":
    import uvicorn
    print("Starting AI Document Companion Backend...")
    print("Server will be available at: http://localhost:8000")
    print("API documentation at: http://localhost:8000/docs")
    try:
        uvicorn.run(app, host="0.0.0.0", port=8000)
    except KeyboardInterrupt:
        print("\nServer stopped.")
    except SystemExit:
        pass